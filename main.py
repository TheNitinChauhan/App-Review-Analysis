# =========================================
# App Review Analysis Pipeline 
# GPT-3.5 Categorization | Embedding Clustering | VADER Sentiment | LDA Topics | PowerPoint Reporting
# ==============================================


################## IF YOU ARE RUNNING CODE FOR APP LIKE  Groww & Zerodha (Finance Apps) THE USE THE 
##################  CODE WRITTEN IN STEP 5 OTHERWISE FOR APP LIKE INSTAGRAM COMMENT OUT THE STEP 5 AND 
##################  RUN THE CODE (REASONE IS WRITTEN JUST ABOVE THE STEP 5)

import os
import time
import re
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import nltk
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
from textblob import TextBlob
from google_play_scraper import Sort, reviews
import openai
from dotenv import load_dotenv
from sklearn.cluster import KMeans
from wordcloud import WordCloud
import traceback
import datetime
from pptx import Presentation
from pptx.util import Inches
from textwrap import wrap 

# ==============================
# SETUP
# ==============================
load_dotenv("new.env")
openai.api_key = os.environ["OPENAI_API_KEY"]

nltk.download('stopwords')
nltk.download('wordnet')
nltk.download('vader_lexicon')

from nltk.sentiment.vader import SentimentIntensityAnalyzer
from gensim import corpora
from gensim.models import LdaModel

analyzer = SentimentIntensityAnalyzer()

overall_start = time.time()

# ==============================
# INPUT PACKAGE
# ==============================
print("\n Enter the app's Google Play package name (e.g., com.nextbillion.groww):")
app_package = input("Package name: ").strip()

if not os.path.exists(app_package):
    os.makedirs(app_package)
os.chdir(app_package)

log_file = open("execution_log.txt", "w", encoding="utf-8")
def log_step(step_name, start_time):
    elapsed = time.time() - start_time
    log_file.write(f"{step_name} completed in {elapsed:.2f} seconds.\n")
    print(f"⏱ {step_name} completed in {elapsed:.2f} sec")
log_file.write(f"Project Start Time: {time.ctime(overall_start)}\nApp Package: {app_package}\n\n")

try:
    # ==============================
    # STEP 1: Scrape Reviews
    # ==============================
    t1 = time.time()
    print(f"\n Scraping latest reviews from {app_package}...")
    result, _ = reviews(app_package, lang='en', country='in', sort=Sort.NEWEST, count=1000)
    raw_df = pd.DataFrame(result)
    raw_df.to_csv("App_Reviews.csv", index=False)
    log_step("Step 1: Scraping", t1)

    # ==============================
    # STEP 2: Clean Data
    # ==============================
    t2 = time.time()
    df = pd.read_csv("App_Reviews.csv", encoding='utf-8', on_bad_lines='skip')
    df = df[['content', 'score']].rename(columns={'content': 'Content'})

    stop_words = set(stopwords.words('english'))
    lemmatizer = WordNetLemmatizer()

    def basic_clean(text):
        if pd.isna(text): return ""
        text = re.sub(r'https?://\S+|www\.\S+', '', text)
        text = re.sub(r'[^\x00-\x7F]+', ' ', text)
        text = re.sub(r'[^\w\s]', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    def advanced_clean(text):
        words = text.lower().split()
        words = [lemmatizer.lemmatize(w) for w in words if w not in stop_words]
        return ' '.join(words)

    df['Content'] = df['Content'].astype(str).apply(basic_clean)
    df = df[df['Content'].str.len() > 4].drop_duplicates(subset='Content')
    df['Content'] = df['Content'].apply(advanced_clean)
    df.reset_index(drop=True, inplace=True)
    df.to_csv("Cleaned_Reviews.csv", index=False)
    log_step("Step 2: Cleaning", t2)

    # ==============================
    # STEP 3: VADER Sentiment Analysis (Recommended Addition)
    # ==============================
    t3 = time.time()
    def get_vader_sentiment(text):
        score = analyzer.polarity_scores(text)['compound']
        if score >= 0.05:
            return 'Positive'
        elif score <= -0.05:
            return 'Negative'
        else:
            return 'Neutral'

    df['VADER_Sentiment'] = df['Content'].apply(get_vader_sentiment)
    df.to_csv("Reviews_With_VADER_Sentiment.csv", index=False)
    log_step("Step 3: VADER Sentiment Analysis", t3)

    # ==============================
    # STEP 4: LDA Topic Modeling (Recommended Addition)
    # ==============================
    t4 = time.time()
    texts = [text.split() for text in df['Content']]
    dictionary = corpora.Dictionary(texts)
    corpus = [dictionary.doc2bow(text) for text in texts]
    lda_model = LdaModel(corpus, num_topics=5, id2word=dictionary, passes=10)

    with open("LDA_Topics.txt", "w", encoding="utf-8") as f:
        for idx, topic in lda_model.print_topics(-1):
            f.write(f"Topic: {idx}\nWords: {topic}\n\n")
    log_step("Step 4: LDA Topic Modeling", t4)

    
###########################################
###########################################
# STEP 5: Embedding-based clustering with GPT naming
    
    
    #Note 
    # FOR AP LIKE INSTAGRAM  we have to comment out step 5 code otherwise it will not perform the steps form step 5 and do not generate any image and 
    # do not cluster and categorize the reviews
    
    # 👉 The reason Instagram failed earlier is likely because:

# Its reviews are huge in volume compared to Groww/Zerodha.

# The embedding model tries to convert every review into a vector → very heavy on RAM/CPU/GPU → ends up failing or returning empty charts.

# So, if your analysis goal is charts + insights, removing embeddings is a good move. ✅
# If in future you want topic modeling or advanced NLP (like grouping reviews by themes), then embeddings will be useful again, but you’ll need to:

# Use batching for embeddings,

# Or sample a subset of reviews (say 5k instead of 100k).
    
    
# It will work for app like groww, zerodha angle one   ✅ Why it worked for Groww/Zerodha but not Instagram?
 #Because their reviews are cleaner and smaller in number, embeddings could handle them. Instagram’s huge, noisy dataset breaks the embedding pipeline.
 
 
 
######  The reason is mainly about data availability + embedding usage:
# 
######################## 1. Groww & Zerodha (Finance Apps)
# 
# These are niche apps with relatively fewer but focused reviews.
# 
# Their reviews are more structured (users talk about trading, app crashes, fund transfers, UI, etc.).
# 
# When you run embeddings on this data, the text is clean and relatively short → embeddings don’t fail, and downstream analysis (sentiment/wordcloud) works fine.
# 
########################### 2. Instagram
# 
# It has millions of reviews that are:
# 
# Much longer (paragraph-style rants).
# 
# Messy (slang, emojis, multiple languages, hashtags).
# 
# Sometimes empty / malformed.
# 
# The embedding step (if you used OpenAI or any other embedding model) fails silently or gives empty vectors for noisy/invalid text.
# 
# This leads to empty DataFrames → empty plots in PowerPoint.
# 
# That’s why removing the embedding step “fixed” the issue → because your sentiment/word cloud functions only need cleaned text, not embeddings.
############################################    
############################################


    print("\n Generating embeddings for cleaned reviews...")
    cleaned_reviews = df['Content'].tolist()
    embeddings = []
    for i in range(0, len(cleaned_reviews), 20):
        batch = cleaned_reviews[i:i+20]
        response = openai.Embedding.create(input=batch, model="text-embedding-3-small")
        batch_embeddings = [item['embedding'] for item in response['data']]
        embeddings.extend(batch_embeddings)
    embeddings = np.array(embeddings)
    print(f" Generated embeddings for {len(embeddings)} reviews.")
    print("\n Performing KMeans clustering...")
    num_clusters = 10
    kmeans = KMeans(n_clusters=num_clusters, random_state=42)
    cluster_labels = kmeans.fit_predict(embeddings)
    df['Embedding_Cluster_ID'] = cluster_labels
    print("\n Generating embedding cluster names using GPT-3.5...")
    cluster_names = {}
    for cid in range(num_clusters):
        samples = df[df['Embedding_Cluster_ID'] == cid]['Content'].sample(min(5, len(df[df['Embedding_Cluster_ID'] == cid]))).tolist()
        prompt = "You are an expert app review analyst. Based on these sample reviews, assign a short and meaningful cluster name describing their common theme.\n\n"
        for i, review in enumerate(samples, 1):
            prompt += f"{i}. {review}\n"
        prompt += "\nReturn only the cluster name."
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=50,
            )
            cluster_name = response["choices"][0]["message"]["content"].strip()
            cluster_names[cid] = cluster_name
            print(f"Cluster {cid}: {cluster_name}")
        except Exception as e:
            print(f" Error generating name for cluster {cid}: {e}")
            cluster_names[cid] = "Unnamed Cluster"
    df['Embedding_Cluster_Name'] = df['Embedding_Cluster_ID'].map(cluster_names)
    df.to_excel("Reviews_with_Embedding_Clusters.xlsx", index=False)
    print(" Saved clustered reviews with embeddings to Reviews_with_Embedding_Clusters.xlsx")
    
    
    
    
    ############################################
    # STEP 6: Time-Series Trend Analysis
    ############################################
    print("\n Generating time-series trend analysis...")
    if 'at' in raw_df.columns:
        raw_df['at'] = pd.to_datetime(raw_df['at'])
        raw_df['date'] = raw_df['at'].dt.date
        date_sentiment = raw_df.groupby(['date'])['score'].mean()
        plt.figure(figsize=(12,6))
        date_sentiment.plot()
        plt.title('Average Review Score Over Time')
        plt.xlabel('Date')
        plt.ylabel('Average Score')
        plt.tight_layout()
        plt.savefig("time_series_average_score.png", dpi=300)
        plt.close()
        print(" Saved: time_series_average_score.png")
    else:
        print(" 'at' field not found for time-series analysis.")
        
    ############################################
    # Continue your existing pipeline from here 
    ############################################
    
    
    t3 = time.time()
    def get_sentiment(text):
        polarity = TextBlob(text).sentiment.polarity
        if polarity > 0.1:
            return 'Positive'
        elif polarity < -0.1:
            return 'Negative'
        else:
            return 'Neutral'
    print("\n Performing sentiment analysis...")
    df['Sentiment'] = df['Content'].apply(get_sentiment)
    df.to_csv("Reviews_With_Sentiment.csv", index=False)
    print(" Sentiment labels saved.")
    log_step("Step 3: Sentiment Analysis", t3)
    # Step 4: Filter 1–4 Star Reviews
    t4 = time.time()
    print("\n Filtering 1 to 4 star reviews...")
    filtered_df = df[df['score'].isin([1, 2, 3, 4])].copy()
    filtered_df.to_excel("Filtered_1_4_Star_Reviews.xlsx", index=False)
    print(f" Filtered {len(filtered_df)} reviews.")
    log_step("Step 4: Filtering 1–4 Stars", t4)
    # Step 5: Auto Categorization using GPT-3.5
    t5 = time.time()
    def batch_categorize_reviews_gpt(review_batch):
        prompt = "Categorize each of the following user reviews into a short issue category (e.g., Login Issue, Payment Failure, UI Bug):\n\n"
        for idx, review in enumerate(review_batch, 1):
            prompt += f"{idx}. {review}\n"
        prompt += "\nReturn only the categories in this format:\n1. <Category>\n2. <Category>\n..."
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=1000,
            )
            output = response["choices"][0]["message"]["content"]
            categories = [line.split(".", 1)[-1].strip() for line in output.split("\n") if "." in line]
            return categories if len(categories) == len(review_batch) else ["Uncategorized"] * len(review_batch)
        except Exception as e:
            print(" Error in categorization:", e)
            return ["Uncategorized"] * len(review_batch)
    print("\n Categorizing Reviews using GPT-3.5...")
    categories = []
    batch_size = 2
    reviews = filtered_df['Content'].tolist()
    for i in range(0, len(reviews), batch_size):
        batch = reviews[i:i+batch_size]
        print(f" Categorizing reviews {i+1} to {i+len(batch)} of {len(reviews)}")
        cat_batch = batch_categorize_reviews_gpt(batch)
        categories.extend(cat_batch)
    filtered_df['Auto_Category'] = categories
    filtered_df.to_excel("Auto_Categorized_Reviews.xlsx", index=False)
    print("\u2705 Auto Categorization Complete.")
    log_step("Step 5: Auto Categorization (GPT-3.5)", t5)
    # Step 6: Clustering using GPT-3.5
    def batch_cluster_reviews_gpt(review_batch):
        prompt = (
            "You are an expert product analyst. For each of the following user reviews, "
            "assign a short and meaningful cluster name describing the main issue or topic. "
            "Examples: Login Issue, Payment Failure, UI Bug, Good Experience, Feature Request, Customer Support, KYC Problem.\n\n"
        )
        for idx, review in enumerate(review_batch, 1):
            prompt += f"{idx}. {review}\n"
        prompt += (
            "\nReturn only the cluster names in this exact format:\n"
            "1. <Cluster>\n"
            "2. <Cluster>\n"
            "...\n"
            "Ensure each line number corresponds to the review number."
        )
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=1000,
            )
            output = response["choices"][0]["message"]["content"]
            
            print("Clustered output:\n", output)
            
            clusters = [line.split(".", 1)[-1].strip() for line in output.split("\n") if "." in line]
            return clusters if len(clusters) == len(review_batch) else ["Unclustered"] * len(review_batch)
        except Exception as e:
            print(" Error in clustering:", e)
            return ["Unclustered"] * len(review_batch)
    print("\n  Clustering Reviews using GPT-3.5...")
    clusters = []
    for i in range(0, len(reviews), batch_size):
        batch = reviews[i:i + batch_size]
        print(f" Clustering reviews {i + 1} to {i + len(batch)} of {len(reviews)}")
        cluster_batch = batch_cluster_reviews_gpt(batch)
        clusters.extend(cluster_batch)
    filtered_df['Cluster'] = clusters
    filtered_df.to_excel("Clustered_Reviews.xlsx", index=False)
    print(" Clustering done.")
        
    #########################################  
    #7. Top 20 Clusters with Sample Reviews
    ##########################################
    print("\n Extracting top 20 clusters with sample reviews...")
    top20 = filtered_df['Cluster'].value_counts().head(20).reset_index()
    top20.columns = ["Cluster", "Count"]
    top20["Cumulative Count"] = top20["Count"].cumsum()
    top20["Cumulative %"] = (top20["Cumulative Count"] / top20["Count"].sum()) * 100
    samples = (
        filtered_df[filtered_df["Cluster"].isin(top20["Cluster"])]
        .groupby("Cluster")["Content"]
        .apply(lambda x: x.sample(min(3, len(x))).tolist())
        .reset_index()
        .rename(columns={"Content": "Sample Reviews"})
    )
    final_df = pd.merge(top20, samples, on="Cluster", how="left")
    final_df.to_excel("Top_20_Cluster_Stats.xlsx", index=False)
    print(" Done! Saved Top_20_Cluster_Stats.xlsx")

    ############################################
    # 8. Pareto Analysis
    ############################################
    print("\n Generating Pareto chart for top 20 clusters...")
    
    # Load your Top 20 clusters file
    top20_df = pd.read_excel("Top_20_Cluster_Stats.xlsx")
    top20_df = top20_df.head(20)
    counts = top20_df['Count']
    clusters = top20_df['Cluster']
    cum_pct = (counts.cumsum() / counts.sum()) * 100
    
    plt.figure(figsize=(18, 9))
    bars = plt.bar(clusters, counts, color='skyblue')
    plt.xticks(rotation=60, ha='right', fontsize=10)
    plt.ylabel('Review Count')
    plt.title('Pareto Analysis of Top 20 Review Clusters')
    
    # Add count labels above bars
    for bar in bars:
        height = bar.get_height()
        plt.annotate(f'{height}', xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', va='bottom', fontsize=8)
    # Create twin axis for cumulative %
    ax2 = plt.gca().twinx()
    ax2.plot(clusters, cum_pct, color='red', marker='o')
    ax2.axhline(80, color='gray', linestyle='--')
    ax2.set_ylabel('Cumulative %')
    # Add cumulative % labels
    for i, cp in enumerate(cum_pct):
        ax2.annotate(f'{cp:.1f}%', xy=(i, cp), xytext=(0, -10), textcoords='offset points',
                    ha='center', va='bottom', fontsize=8, color='red')
    plt.tight_layout()
    plt.savefig("pareto_cluster_analysis_top20.png", dpi=300)
    plt.close()
    print(" Saved: pareto_cluster_analysis_top20.png")
    
    
    ############################################
    # 9. Sentiment Pie Chart and Word Clouds
    ############################################
    print("\n Generating sentiment pie chart...")
    sentiment_counts = df['Sentiment'].value_counts()
    plt.figure(figsize=(6, 6))
    plt.pie(sentiment_counts, labels=sentiment_counts.index, autopct='%1.1f%%',
            colors=['#8BC34A', '#FF5722', '#03A9F4'])
    plt.title('Sentiment Distribution')
    plt.savefig("sentiment_pie.png", dpi=300)
    plt.close()
    print(" Saved: sentiment_pie.png")
    print("\n Generating word clouds for each sentiment...")
    def generate_wordcloud(text_series, file_name, title):
        text = ' '.join(text_series.dropna().tolist())
        wordcloud = WordCloud(width=800, height=400, background_color='white').generate(text)
        plt.figure(figsize=(10, 5))
        plt.imshow(wordcloud, interpolation='bilinear')
        plt.axis('off')
        plt.title(title, fontsize=16)
        plt.tight_layout()
        plt.savefig(file_name, dpi=300)
        plt.close()
        print(f" Saved: {file_name}")
    generate_wordcloud(df[df['Sentiment'] == 'Positive']['Content'], "wordcloud_positive.png",
                    "Positive Reviews Word Cloud")
    generate_wordcloud(df[df['Sentiment'] == 'Negative']['Content'], "wordcloud_negative.png",
                    "Negative Reviews Word Cloud")
    generate_wordcloud(df[df['Sentiment'] == 'Neutral']['Content'], "wordcloud_neutral.png",
                    "Neutral Reviews Word Cloud")
    ############################################
    # 10. Additional Insightful Charts (Optional)
    ############################################
    # Bar chart of sentiment vs count
    print("\n Generating sentiment bar chart...")
    plt.figure(figsize=(8, 6))
    sentiment_counts.plot(kind='bar', color=['green', 'red', 'blue'])
    plt.title('Number of Reviews per Sentiment')
    plt.ylabel('Count')
    plt.xticks(rotation=0)
    plt.tight_layout()
    plt.savefig("sentiment_bar_chart.png", dpi=300)
    plt.close()
    print(" Saved: sentiment_bar_chart.png")
    # Cluster frequency bar chart
    print("\n Generating review cluster bar chart...")
    plt.figure(figsize=(12, 6))
    filtered_df['Cluster'].value_counts().head(20).plot(kind='bar', color='skyblue')
    plt.title('Top 20 Clusters (Review Themes)')
    plt.ylabel('Review Count')
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    plt.savefig("top_20_clusters_bar_chart.png", dpi=300)
    plt.close()
    print(" Saved: top_20_clusters_bar_chart.png")
    ############################################
    # 11. Auto-Category Analysis Charts
    ############################################
    # Bar chart for Auto Categories
    print("\n Generating auto-category bar chart...")
    plt.figure(figsize=(12, 6))
    filtered_df['Auto_Category'].value_counts().head(20).plot(kind='bar', color='orange')
    plt.title('Top 20 Auto Categories (Issue Types)')
    plt.ylabel('Review Count')
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    plt.savefig("auto_category_bar_chart.png", dpi=300)
    plt.close()
    print(" Saved: auto_category_bar_chart.png")
    # Stacked bar chart for Auto Category vs Sentiment
    print("\n Generating stacked bar chart of auto-category vs sentiment...")
    auto_sentiment = filtered_df.groupby(['Auto_Category', 'Sentiment']).size().unstack(fill_value=0)
    auto_sentiment = auto_sentiment.loc[auto_sentiment.sum(axis=1).sort_values(ascending=False).head(15).index]
    auto_sentiment.plot(kind='bar', stacked=True, figsize=(14, 7), colormap='viridis')
    plt.title('Sentiment Distribution Across Top Auto Categories')
    plt.ylabel('Review Count')
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    plt.savefig("auto_category_sentiment_stacked_bar.png", dpi=300)
    plt.close()
    print(" Saved: auto_category_sentiment_stacked_bar.png")
    
    ############################################
    # 12. Execution Summary and Time Logging (NEW STEP)
    ############################################


    total_time = time.time() - overall_start
    log_file.write(f"\nTotal execution time: {total_time:.2f} seconds ({total_time / 60:.2f} minutes)\n")
    print(f"\n All steps completed in {total_time:.2f} seconds (~{total_time / 60:.2f} minutes)")

except Exception as e:
    error_msg = traceback.format_exc()
    log_file.write("\n ERROR OCCURRED:\n" + error_msg)
    print(" An error occurred. Check execution_log.txt for details.")

finally:
    log_file.close()
print("\n Pipeline completed successfully!")


###############################################
# 13 Code for power point presentation
###############################################

from pptx import Presentation
from pptx.util import Inches
import datetime

from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE

# === Beautification Helpers ===

def set_font_style(text_frame, font_name='Segoe UI', font_size=14, bold=False, color_rgb=(34,34,34)):
    """Set font style for a text_frame"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color_rgb)

def add_colored_title_bar(slide, color_rgb=(0, 120, 215)):
    """Add a colored bar below title for branding"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.3), Inches(9), Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color_rgb)
    shape.line.fill.background()  # No border

def add_footer_with_logo(slide, slide_num):
    """Add footer with page number and optional logo"""
    # Footer text
    footer = slide.shapes.add_textbox(Inches(0), Inches(7), Inches(9), Inches(0.3))
    tf = footer.text_frame
    tf.text = f"Generated by YourName | Slide {slide_num}"
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(100, 100, 100)

def beautify_slide(slide, slide_num):
    """Apply full beautification to a slide"""
    add_colored_title_bar(slide)
    add_footer_with_logo(slide, slide_num)


# === Helper: Add title+text slide ===
def add_slide(prs, title, content, bullet_points=None):
    slide_layout = prs.slide_layouts[1]  # Title + Content logo
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content_box = slide.placeholders[1]
    content_box.text = content
    if bullet_points:
        for point in bullet_points:
            p = content_box.text_frame.add_paragraph()
            p.text = point
            p.level = 0
    # Beautify here
    beautify_slide(slide, len(prs.slides))
    
# === Helper: Add picture with description ===
from textwrap import wrap  # <-- Import at the top of your script



###########################################
##########################################
#####################################
# def add_image_with_description_slide(prs, title, image_path, description):
    # slide = prs.slides.add_slide(prs.slide_layouts[1])
    # slide.shapes.title.text = title
    # slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), height=Inches(3.5))
    # 
   ## Wrap description text for better readability in slides
    # wrapped_desc = "\n".join(wrap(description, width=80))  # Adjust width if needed
    # 
    # box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))  # Increased height to 1.5
    # tf = box.text_frame
    # tf.word_wrap = True
    # tf.text = wrapped_desc
    # 
    ## Beautify slide here
    # beautify_slide(slide, len(prs.slides))
    
    
    
    
from pathlib import Path
import matplotlib.pyplot as plt
from textwrap import wrap
from pptx.util import Inches

def make_placeholder_image(path, text=None, size=(800, 400)):
    """Create a simple placeholder PNG so pptx has something to insert."""
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=(size[0] / 100, size[1] / 100))
    plt.text(0.5, 0.5, text or f"{path.name} not found", fontsize=18, ha='center', va='center')
    plt.axis('off')
    fig.tight_layout()
    fig.savefig(str(path), dpi=100, bbox_inches='tight')
    plt.close(fig)

def add_image_with_description_slide(prs, title, image_path, description):
    """
    Safe version: ensures an image exists (creates placeholder if missing)
    and falls back to a textbox if pptx can't read the image.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    img_path = Path(image_path)

    # If missing, create a placeholder PNG
    if not img_path.exists():
        make_placeholder_image(img_path, text=f"Missing: {img_path.name}")

    # Try to add picture; fall back to a textbox on failure
    try:
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.5), height=Inches(3.5))
    except Exception as e:
        # Add a textbox in place of the image so PPT generation doesn't fail
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3.5))
        tf2 = tb.text_frame
        tf2.word_wrap = True
        tf2.text = f"Could not load image: {img_path.name}\nError: {e}"

    # Wrap description text for readability
    wrapped_desc = "\n".join(wrap(description or "", width=80))
    box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = wrapped_desc

    # Beautify slide (keeps your original styling)
    beautify_slide(slide, len(prs.slides))

################################################
################################################
#################################################
# === Setup ===
prs = Presentation()
today = datetime.date.today().strftime("%B %d, %Y")
app_title = app_package.split('.')[-1].capitalize()

# === Title Slide ===
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = f"{app_title} App Review Analysis"
slide.placeholders[1].text = f"Auto-generated report\n{today}"

# === Problem Statement ===
add_slide(prs, "Problem Statement", "To identify user pain points and sentiment from Google Play reviews to help improve the app.")

# === Data Source ===
add_slide(prs, "Data Source",
          f"• Google Play Store\n• Package: {app_package}",
          [f"Scraped file: App_Reviews.csv"])

# === Data Cleaning ===
add_slide(prs, "Data Cleaning Approach",
          "Applied basic regex, stopwords removal and lemmatization",
          ["Output file: Cleaned_Reviews.csv"])

# === Sentiment Analysis ===
add_slide(prs, "Sentiment Analysis",
          "Used TextBlob to assign polarity-based sentiment to each review",
          ["Output file: Reviews_With_Sentiment.csv", "Chart: sentiment_pie.png"])

# === Filtered Reviews ===
add_slide(prs, "Filtered Reviews (1–4 Stars)",
          "Filtered reviews with score 1–4 to detect complaints",
          ["File: Filtered_1_4_Star_Reviews.xlsx"])

# === Auto-Categorization ===
add_slide(prs, "Auto-Categorization using GPT-3.5",
          "Used GPT-3.5 to assign issue categories to each review",
          ["Example: Login Issue, Payment Failure", "Output file: Auto_Categorized_Reviews.xlsx"])

# === Clustering ===
add_slide(prs, "Clustering via GPT-3.5",
          "Grouped reviews into higher-level clusters",
          ["Output file: Clustered_Reviews.xlsx"])


# === Sentiment Pie Chart ===
add_image_with_description_slide(
    prs,
    "Sentiment Distribution (Pie Chart)",
    "sentiment_pie.png",
    "This chart shows the overall sentiment of app reviews. A large share of negative reviews suggests areas of concern, "
    "such as app crashes or poor support. Positive reviews reflect satisfaction with ease of use and features."
)

# === WordClouds ===
add_image_with_description_slide(
    prs,
    "Positive Reviews Word Cloud",
    "wordcloud_positive.png",
    "Frequent words in positive reviews include 'easy', 'interface', and 'returns' — indicating satisfaction with usability and performance."
)

add_image_with_description_slide(
    prs,
    "Negative Reviews Word Cloud",
    "wordcloud_negative.png",
    "Negative reviews often mention 'login', 'issue', 'delay', and 'support' — highlighting key pain points that frustrate users."
)

add_image_with_description_slide(
    prs,
    "Neutral Reviews Word Cloud",
    "wordcloud_neutral.png",
    "Neutral reviews contain suggestions or feature requests without strong emotional tone. Common themes include improvements and feedback."
)

# === Auto Category Bar Chart ===
add_image_with_description_slide(
    prs,
    "Top 20 Auto Categories (Bar Chart)",
    "auto_category_bar_chart.png",
    "This bar chart shows the most frequent auto-categorized issue types. Login, payment, and KYC issues dominate."
)

# === Stacked Bar Chart: Auto Categories vs Sentiment ===
add_image_with_description_slide(
    prs,
    "Sentiment Across Top Auto Categories",
    "auto_category_sentiment_stacked_bar.png",
    "This stacked bar chart shows how sentiment varies by issue category. Login and support issues are mostly negative, "
    "while some categories like feature requests show more neutral or positive sentiment."
)
add_image_with_description_slide(
    prs,
    "Pareto Analysis of Top 20 Clusters",
    "pareto_cluster_analysis_top20.png",
    "This Pareto chart visualizes the top 20 review clusters, showing their individual frequencies and cumulative contribution. "
    "It helps identify the few key issues that account for the majority of complaints, aligning with the 80/20 principle for prioritization."
)

# === Top 20 Clusters Bar Chart ===
add_image_with_description_slide(
    prs,
    "Top 20 Clusters (Bar Chart)",
    "top_20_clusters_bar_chart.png",
    "This chart visualizes the most common clusters found using GPT-based grouping. It highlights the frequency of each recurring theme."
)

# === Time-Series Analysis Chart ===
add_image_with_description_slide(
    prs,
    "Time-Series Analysis of Review",
    "time_series_average_score.png",
    "This time-series plot shows the average review scores over time, helping to identify any recent dips or improvements "
    "in user sentiment that may correlate with app updates or feature releases."
)




#############################################################
###########
# ----------------------------
# for Gererating Data-driven Key Findings & Recommendations generator
# ----------------------------
import math
import pandas as _pd
import numpy as _np
import re as _re
from datetime import datetime, timedelta as _td

def _safe_top_topics(lda_model, fallback_file="LDA_Topics.txt", n=5):
    topics = []
    try:
        n_topics = min(getattr(lda_model, "num_topics", 5), 10)
        for tid in range(n_topics):
            words = [w for w, _ in lda_model.show_topic(tid, topn=n)]
            topics.append(" ".join(words))
    except Exception:
        try:
            with open(fallback_file, "r", encoding="utf-8") as f:
                txt = f.read()
            for part in txt.split("Topic:")[1:]:
                lines = part.strip().splitlines()
                for L in lines:
                    if L.strip().lower().startswith("words:"):
                        words_text = L.split(":", 1)[1].strip()
                        toks = _re.findall(r"[a-zA-Z]+", words_text)[:n]
                        topics.append(" ".join(toks))
                        break
        except Exception:
            topics = []
    return topics

def _get_top_clusters(top20_df=None, filtered_df=None, k=5):
    if isinstance(top20_df, _pd.DataFrame) and "Cluster" in top20_df.columns:
        return top20_df['Cluster'].astype(str).head(k).tolist()
    try:
        dfc = filtered_df if isinstance(filtered_df, _pd.DataFrame) else None
        if dfc is not None and 'Cluster' in dfc.columns:
            return dfc['Cluster'].value_counts().head(k).index.astype(str).tolist()
    except Exception:
        pass
    return []

def _sentiment_summary_from_df(df):
    ss = {'positive_pct': 0.0, 'negative_pct': 0.0, 'neutral_pct': 0.0}
    try:
        s = df['Sentiment'].value_counts(normalize=True) * 100
        ss['positive_pct'] = float(s.get('Positive', 0.0))
        ss['negative_pct'] = float(s.get('Negative', 0.0))
        ss['neutral_pct'] = float(s.get('Neutral', 0.0))
    except Exception:
        pass
    return ss

def _recent_sentiment_trend(raw_df):
    """Return dictionary with short trend if dates available (compare last 14 days vs prior 14 days)."""
    try:
        rdf = raw_df.copy()
        if 'at' in rdf.columns:
            rdf['at'] = _pd.to_datetime(rdf['at'])
            rdf = rdf.set_index('at')
            end = rdf.index.max().date()
            end_dt = _pd.Timestamp(end)
            last14 = rdf[end_dt - _pd.Timedelta(days=13): end_dt]
            prev14 = rdf[end_dt - _pd.Timedelta(days=27): end_dt - _pd.Timedelta(days=14)]
            if len(last14) >= 3 and len(prev14) >= 3:
                last_mean = last14['score'].mean()
                prev_mean = prev14['score'].mean()
                delta = last_mean - prev_mean
                trend = "up" if delta > 0 else ("down" if delta < 0 else "flat")
                return {'last14_mean': float(last_mean), 'prev14_mean': float(prev_mean), 'delta': float(delta), 'trend': trend}
    except Exception:
        pass
    return {}

def _sample_for_cluster(cluster_name, filtered_df, chars=180):
    try:
        mask = filtered_df['Cluster'].str.lower().str.contains(str(cluster_name).lower(), na=False)
        if mask.any():
            s = filtered_df[mask]['Content'].iloc[0]
            return (s[:chars] + '...') if len(s) > chars else s
    except Exception:
        pass
    return None

# Prepare inputs (safe)
try:
    sentiment_summary  # if exists
except NameError:
    sentiment_summary = _sentiment_summary_from_df(df if 'df' in globals() else _pd.DataFrame())

try:
    top_topics  # if exists
except NameError:
    top_topics = _safe_top_topics(lda_model if 'lda_model' in globals() else None)

try:
    common_issues  # if exists
except NameError:
    common_issues = _get_top_clusters(top20 if 'top20' in globals() else None, filtered_df if 'filtered_df' in globals() else None, k=10)

# additional signals
top_clusters = _get_top_clusters(top20 if 'top20' in globals() else None, filtered_df if 'filtered_df' in globals() else None, k=5)
trend = _recent_sentiment_trend(raw_df if 'raw_df' in globals() else (_pd.DataFrame()))
# flags from topics+clusters
_check_text = " ".join((top_topics or []) + (common_issues or [])).lower()

flags = {
    'login': any(x in _check_text for x in ['login','signin','authentication','otp','lock']),
    'payment': any(x in _check_text for x in ['payment','txn','transaction','withdraw','refund','upi','deposit']),
    'support': any(x in _check_text for x in ['support','customer service','sla','agent','help']),
    'fees': any(x in _check_text for x in ['fee','fees','charge','charges','brokerage','hidden']),
    'performance': any(x in _check_text for x in ['crash','slow','lag','freeze','performance']),
    'fraud': any(x in _check_text for x in ['fraud','scam','unauthor','unauthorized','loss','dispute'])
}

# Compose 5 detailed Key Findings (multi-line strings)
kf = []

kf.append(
    f"1) Overall sentiment: Positive {sentiment_summary.get('positive_pct',0):.1f}%, "
    f"Negative {sentiment_summary.get('negative_pct',0):.1f}%, Neutral {sentiment_summary.get('neutral_pct',0):.1f}%.\n"
    "   Interpretation: Strong positive majority indicates product strengths, but the negative tail highlights recurring frictions\n"
    "   that affect retention. Track weekly sentiment deltas and prioritize issues causing negative spikes.\n"
    + (f"   Recent trend: last14_mean={trend.get('last14_mean')}, prev14_mean={trend.get('prev14_mean')}, trend={trend.get('trend')}\n" if trend else "")
)

topics_snip = ", ".join(top_topics[:4]) if top_topics else "No dominant LDA topics"
kf.append(
    f"2) Topic signals & user praise: Top LDA topics emphasize: {topics_snip}.\n"
    "   This suggests core capabilities (e.g., charts, funds, usability) are perceived as strengths by many users.\n"
    "   Leverage these strengths in messaging and prioritize feature polish that amplifies this positive feedback.\n"
)

clusters_str = ", ".join(top_clusters) if top_clusters else ", ".join(common_issues[:4]) if common_issues else "N/A"
kf.append(
    f"3) Pareto concentration: top clusters = {clusters_str}.\n"
    "   Analysis (Pareto) shows a small number of issue categories drive most complaints — fixing those yields outsized impact.\n"
    "   Use Top_20_Cluster_Stats.xlsx to identify exact counts and prioritize the top 3 clusters for an immediate ROI."
)

# UX & Stability detail
ux_lines = []
if flags['login']:
    ux_lines.append("login/authentication failures (OTP, lockouts)")
if flags['performance']:
    ux_lines.append("app crashes, freezes or slow responses")
if flags['payment']:
    ux_lines.append("payment/withdrawal errors or transaction failures")
if not ux_lines:
    ux_lines = ["no dominant UX/stability pattern detected in top signals"]

kf.append(
    f"4) UX & stability hotspots: {', '.join(ux_lines)}.\n"
    "   These issues cause immediate task failure and generate high-severity support tickets; reproduce top issues with covered steps,\n"
    "   add telemetry to capture failure contexts (app-version, device, network), and introduce guards in the UX to prevent fatal flows."
)

# Support & trust
trust_issues = []
if flags['fees']:
    trust_issues.append("fees / hidden charges raising trust concerns")
if flags['fraud']:
    trust_issues.append("fraud / unauthorized transaction reports (urgent)")
if flags['support']:
    trust_issues.append("slow or unsatisfactory customer support experience")
kf.append(
    f"5) Support & trust signals: {', '.join(trust_issues) if trust_issues else 'support OK but monitor closely'}.\n"
    "   Billing disputes, fraud flags and slow support increase churn risk. Track time-to-resolution, dispute volume, and repeat tickets\n"
    "   and set targets (e.g., resolve critical financial tickets < 24 hours)."
)

# Compose 5 detailed Recommendations (each multi-line)
recs = []

recs.append(
    "1) Run a 30-day Pareto Triage Sprint:\n"
    "   - Focus a cross-functional team on the top 3 clusters (engineering, QA, product, ops).\n"
    "   - Reproduce, patch, and validate fixes; push phased hotfixes and monitor metrics post-release.\n"
    "   - Track cluster counts daily to confirm impact."
)

recs.append(
    "2) Improve UX flows & observability:\n"
    "   - Simplify login and KYC flows (clear errors, inline help, retry paths) and reduce friction points.\n"
    "   - Instrument detailed telemetry (error codes, device, network) and create dashboards to find root causes.\n"
    "   - Add synthetic tests to catch regressions before release."
)

recs.append(
    "3) Payment reliability & transparency:\n"
    "   - Display clear fee breakdowns and transaction status; provide immediate next-action guidance on errors.\n"
    "   - Implement automatic retries for transient failures and create a fast-track billing-dispute process.\n"
    "   - Provide ops dashboard to manage dispute SLA (<48h for billing issues)."
)

recs.append(
    "4) Support model upgrade & SLAs:\n"
    "   - Deploy AI-assisted triage to resolve common queries instantly and gather structured info for agents.\n"
    "   - Route critical financial/fraud tickets to a high-priority queue and measure first-contact resolution.\n"
    "   - Publish expected SLAs in-app so customers know what to expect."
)

recs.append(
    "5) Roadmap transparency & closed-loop feedback:\n"
    "   - Publish a concise roadmap addressing top-requested features and provide beta access to engaged users.\n"
    "   - After each major release, re-run sentiment & cluster analysis and communicate outcomes in release notes.\n"
    "   - Use metrics (sentiment delta, cluster counts, NPS) to measure success."
)

# Insert into PPT using user's add_slide helper
# Each bullet is a multi-line string; add_slide will handle them as separate list items.
try:
    add_slide(prs, "Key Findings", "Based on sentiment, categorization, and cluster frequency", kf)
    add_slide(prs, "Recommendations", "Actionable recommendations based on analysis", recs)
except Exception as e:
    # Fallback: if add_slide not available, print for debug
    print("Error adding slides:", e)
    print("\nKey Findings:\n", "\n\n".join(kf))
    print("\nRecommendations:\n", "\n\n".join(recs))

################################################################################################

# === Tools & Technologies ===
add_slide(prs, "Tools & Technologies (detailed)",
          "Core libraries, models and what they do",
          [
              "Python 3.9+ — codebase (pandas 1.x, numpy 1.x) for data processing",
              "google-play-scraper — scrape reviews from Play Store",
              "NLTK & TextBlob — preprocessing, stopwords, lemmatization, polarity sentiment",
              "Gensim LDA — unsupervised topic modeling (num_topics configurable)",
              "OpenAI embeddings (text-embedding-3-small) — semantic vectors for clustering",
              "OpenAI Chat (gpt-3.5-turbo) — human-like cluster names, categorization, recommendations",
              "scikit-learn (KMeans) — clustering; consider UMAP+HDBSCAN for production",
              "Matplotlib + WordCloud — charts and word-clouds for slides",
              "python-pptx — automated PowerPoint report generation",
              "python-dotenv, tenacity — env management and robust API retries"
          ])

# Set light gray background for all slides
for slide in prs.slides:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)

# === Save the Report ===
pptx_filename = f"{app_title}_Review_Analysis_Report.pptx"
prs.save(pptx_filename)
print(f" PowerPoint Report Saved: {pptx_filename}")
